# app.py
import os
import re
import json
import math

import psycopg2
from db import get_db_conn, init_db
import requests
import tempfile
from datetime import datetime
from flask import Flask, request, jsonify, render_template, send_from_directory, Response, session
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from dotenv import load_dotenv

# your curriculum_final.py must be in same folder
from curriculum_final import (
    extract_text_any, extract_structure_any,
    generate_curriculum_9_months, generate_curriculum_custom, highlight_excel, export_to_word
)

# ------------- CONFIG -------------
load_dotenv()
APP_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FOLDER = os.path.join(APP_ROOT, "outputs")
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app = Flask(__name__)

@app.errorhandler(psycopg2.OperationalError)
def handle_db_operational_error(e):
    return jsonify({
        "error": "Database Service Unavailable",
        "message": "Could not connect to Supabase PostgreSQL database. Please verify connection credentials and network status."
    }), 503
app.secret_key = os.getenv("SECRET_KEY", "dev_default_secret_key_change_in_production")
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['SESSION_COOKIE_SECURE'] = os.getenv("FLASK_ENV") == "production"

# HARDENED PRODUCTION CORS ORIGIN RESTRICTION
# Restrict API access exclusively to trusted frontend domains (e.g. localhost for dev, all *.vercel.app domains for prod)
raw_origins = os.getenv("ALLOWED_ORIGINS") or os.getenv("FRONTEND_URL") or "http://localhost:3000,http://127.0.0.1:3000,http://localhost:5173,https://robolearn.vercel.app"
allowed_origins = [o.strip() for o in raw_origins.split(",") if o.strip()]
# Add regex matcher for all Vercel deployment & preview URLs
allowed_origins.append(re.compile(r"https://.*\.vercel\.app"))

CORS(
    app,
    supports_credentials=True,
    origins=allowed_origins,
    methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["Content-Type", "Authorization", "X-Requested-With"]
)

@app.after_request
def add_cors_headers(response):
    origin = request.headers.get("Origin")
    if origin:
        if origin.endswith(".vercel.app") or "localhost" in origin or "127.0.0.1" in origin:
            response.headers["Access-Control-Allow-Origin"] = origin
            response.headers["Access-Control-Allow-Credentials"] = "true"
            response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, DELETE, OPTIONS"
            response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, X-Requested-With"
    return response



# ------------- DB / MIGRATION HELPERS -------------
# Database connectivity is initialized via db.py using Supabase PostgreSQL connection pool.

# Helper to check for guest / unauthenticated sessions
def is_guest_user(user_id):
    """
    Returns True if user_id represents an unauthenticated/guest session (e.g. 'student', None, 'guest').
    Guest sessions are not persisted to database tables to prevent Foreign Key constraint violations.
    """
    if not user_id:
        return True
    return str(user_id).strip().lower() in ["student", "guest", "anonymous", "null", "none", ""]

def save_memory(user_id, book_text=None, pastpaper_text=None, toc=None):
    if is_guest_user(user_id):
        # Transient guest session: skip DB persistence
        return
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("SELECT book_text, pastpaper_text, toc_json FROM teacher_memory WHERE user_id=%s", (user_id,))
        row = c.fetchone()
        old_book, old_past, old_toc = ("", "", "[]") if not row else (row[0] or "", row[1] or "", row[2] or "[]")
        new_book = book_text if book_text is not None else old_book
        new_past = pastpaper_text if pastpaper_text is not None else old_past
        new_toc = json.dumps(toc, ensure_ascii=False) if toc is not None else old_toc
        c.execute("INSERT INTO teacher_memory (user_id, book_text, pastpaper_text, toc_json) VALUES (%s, %s, %s, %s) ON CONFLICT (user_id) DO UPDATE SET book_text=EXCLUDED.book_text, pastpaper_text=EXCLUDED.pastpaper_text, toc_json=EXCLUDED.toc_json", (user_id, new_book, new_past, new_toc))
        conn.commit(); conn.close()
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [save_memory]: Failed to save memory for user '{user_id}': {e}")

def get_memory(user_id):
    if is_guest_user(user_id):
        # Transient guest session: return empty memory
        return "", "", []
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("SELECT book_text, pastpaper_text, toc_json FROM teacher_memory WHERE user_id=%s", (user_id,))
        row = c.fetchone(); conn.close()
        if not row:
            return "", "", []
        book, past, toc_json = row[0] or "", row[1] or "", row[2] or "[]"
        try:
            toc = json.loads(toc_json)
        except Exception:
            toc = []
        return book, past, toc
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [get_memory]: Failed to fetch memory for user '{user_id}': {e}")
        return "", "", []

def add_message(user_id, role, content):
    if is_guest_user(user_id):
        # Transient guest session: skip message persistence to DB
        return
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("INSERT INTO messages (user_id, role, content, created_at) VALUES (%s, %s, %s, NOW())", (user_id, role, content))
        conn.commit(); conn.close()
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [add_message]: Failed to insert message for user '{user_id}': {e}")

def get_history(user_id, limit=50):
    if is_guest_user(user_id):
        # Transient guest session: return empty chat history
        return []
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("SELECT role, content FROM messages WHERE user_id=%s ORDER BY id ASC LIMIT %s", (user_id, limit))
        rows = c.fetchall(); conn.close()
        return [{"role": r[0], "content": r[1]} for r in rows]
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [get_history]: Failed to fetch chat history for user '{user_id}': {e}")
        return []

def save_quiz(user_id, quiz_obj):
    if is_guest_user(user_id):
        # Transient guest session: skip quiz persistence
        return
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("INSERT INTO quizzes (user_id, quiz_json, created_at) VALUES (%s, %s, NOW())", (user_id, json.dumps(quiz_obj, ensure_ascii=False)))
        conn.commit(); conn.close()
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [save_quiz]: Failed to save quiz for user '{user_id}': {e}")

def get_latest_quiz(user_id):
    if is_guest_user(user_id):
        return None
    try:
        conn = get_db_conn(); c = conn.cursor()
        c.execute("SELECT quiz_json FROM quizzes WHERE user_id=%s ORDER BY id DESC LIMIT 1", (user_id,))
        row = c.fetchone(); conn.close()
        return json.loads(row[0]) if row else None
    except psycopg2.Error as e:
        print(f"DATABASE WARNING [get_latest_quiz]: Failed to fetch latest quiz for user '{user_id}': {e}")
        return None

# ------------- GEMINI API & WEB SEARCH HELPERS -------------
from config import GEMINI_API_KEY, GEMINI_MODEL, call_gemini, stream_gemini

def perform_web_search(query, max_results=3):
    """
    Perform real-time web search using Tavily API.
    Returns tuple: (formatted_text_context, list_of_source_dicts)
    """
    if not query or len(query.strip()) < 3:
        return "", []

    tavily_key = os.getenv("TAVILY_API_KEY", "")
    if not tavily_key:
        print("Tavily API warning: TAVILY_API_KEY environment variable is missing.")
        return "", []

    try:
        from tavily import TavilyClient
        client = TavilyClient(api_key=tavily_key)
        response = client.search(query=query, max_results=max_results, search_depth="basic", timeout=5.0)

        results = response.get("results", []) if isinstance(response, dict) else []
        results_text = []
        sources = []

        for r in results:
            title = r.get("title", "")
            snippet = r.get("content", "") or r.get("snippet", "")
            url = r.get("url", "")
            if snippet or title:
                results_text.append(f"• Title: {title}\n  Snippet: {snippet}\n  URL: {url}")
                sources.append({"title": title, "url": url, "snippet": snippet})

        if results_text:
            return "Web Search Real-Time Context:\n\n" + "\n\n".join(results_text), sources
    except Exception as e:
        print("Tavily Web Search Warning/Error:", e)

    return "", []



# ------------- ADVANCED SEMANTIC RAG & VECTOR SEARCH ENGINE -------------
import numpy as np
_embedding_model = None

def get_embedding_model():
    global _embedding_model
    if _embedding_model is None:
        from sentence_transformers import SentenceTransformer
        _embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    return _embedding_model

def generate_and_store_embeddings(book_id, book_text, conn):
    """Chunk book text into ~400-token pieces, find exact chapter ranges, and store embeddings."""
    model = get_embedding_model()
    
    # 1. Retrieve all chapters for this book
    c = conn.cursor()
    c.execute("SELECT id, title, chapter_number FROM chapters WHERE book_id=%s ORDER BY chapter_number", (book_id,))
    chapter_rows = c.fetchall()
    
    # 2. Estimate ranges in raw book text
    # Map each chapter to a character range [start, end)
    ch_ranges = []
    for i, (ch_id, ch_title, ch_num) in enumerate(chapter_rows):
        # We need to find the start offset of this chapter's title.
        # To avoid matching in Table of Contents at start of book, skip matches under character offset 10000 (if a later one exists)
        title_lower = ch_title.lower()
        offset = book_text.lower().find(title_lower)
        if offset != -1 and offset < 10000:
            next_offset = book_text.lower().find(title_lower, offset + len(ch_title))
            if next_offset != -1:
                offset = next_offset
        
        # Fallbacks if title not found
        if offset == -1:
            ch_pat = f"chapter {ch_num}"
            offset = book_text.lower().find(ch_pat)
            if offset != -1 and offset < 10000:
                next_offset = book_text.lower().find(ch_pat, offset + len(ch_pat))
                if next_offset != -1:
                    offset = next_offset
                    
            if offset == -1:
                ch_pat = f"unit {ch_num}"
                offset = book_text.lower().find(ch_pat)
                if offset != -1 and offset < 10000:
                    next_offset = book_text.lower().find(ch_pat, offset + len(ch_pat))
                    if next_offset != -1:
                        offset = next_offset
        
        ch_ranges.append({
            "id": ch_id,
            "title": ch_title,
            "start": offset if offset != -1 else 0
        })
    
    # Sort by start offset to establish intervals
    ch_ranges.sort(key=lambda x: x["start"])
    for i in range(len(ch_ranges)):
        start = ch_ranges[i]["start"]
        end = ch_ranges[i+1]["start"] if (i + 1) < len(ch_ranges) else len(book_text)
        ch_ranges[i]["end"] = end

    # 3. Chunk text into ~1500 character window paragraphs
    paragraphs = [p.strip() for p in book_text.split('\n\n') if len(p.strip()) > 30]
    chunks = []
    current_chunk = ""
    current_offset = 0
    
    for p in paragraphs:
        if len(current_chunk) + len(p) < 1500:
            current_chunk += "\n\n" + p if current_chunk else p
        else:
            if current_chunk:
                chunks.append((current_chunk, current_offset))
            current_chunk = p
            current_offset = book_text.find(p, current_offset)
    if current_chunk:
        chunks.append((current_chunk, current_offset))
    
    if not chunks:
        return
        
    # Batch encode all chunks
    texts = [c[0] for c in chunks]
    embeddings = model.encode(texts, show_progress_bar=False, normalize_embeddings=True)
    
    # 4. Insert chunks mapped to correct range-based chapter IDs
    for i, (chunk_text, offset) in enumerate(chunks):
        emb_blob = embeddings[i].astype(np.float32).tobytes()
        
        # Match chunk offset to chapter range
        ch_id = chapter_rows[0][0] if chapter_rows else None
        for rng in ch_ranges:
            if rng["start"] <= offset < rng["end"]:
                ch_id = rng["id"]
                break
        
        c.execute("""
            INSERT INTO chunk_embeddings (book_id, chapter_id, chunk_index, chunk_text, char_offset, embedding)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (book_id, ch_id, i, chunk_text, offset, psycopg2.Binary(emb_blob)))

STOP_WORDS = {
    'help', 'me', 'learn', 'learning', 'teach', 'teaching', 'this', 'make', 'mem',
    'please', 'want', 'to', 'can', 'you', 'explain', 'tell', 'about', 'what', 'is',
    'are', 'the', 'a', 'an', 'in', 'on', 'of', 'for', 'with', 'book', 'textbook',
    'how', 'does', 'do', 'why', 'who', 'when', 'where', 'which'
}

# Domain Concept & Synonym Expansion Dictionary
CONCEPT_SYNONYMS = {
    "cpu": ["central processing unit", "processor", "execution", "registers", "alu", "control unit", "clock", "instruction"],
    "memory": ["ram", "cache", "storage", "dram", "sram", "registers", "address space", "hierarchy"],
    "speed": ["performance", "latency", "throughput", "bandwidth", "clock rate", "execution time"],
    "network": ["protocol", "tcp", "ip", "packets", "router", "switch", "bandwidth", "socket"],
    "database": ["sql", "query", "index", "table", "schema", "relational", "transaction", "acid"],
    "code": ["program", "algorithm", "function", "variable", "syntax", "compilation", "execution"],
    "ai": ["artificial intelligence", "machine learning", "neural network", "deep learning", "model", "training"],
    "security": ["encryption", "authentication", "authorization", "hash", "cipher", "vulnerability", "firewall"],
    "os": ["operating system", "kernel", "process", "thread", "scheduling", "virtual memory", "deadlock"],
    "hardware": ["circuit", "transistor", "gate", "bus", "microarchitecture", "silicon", "motherboard"]
}

def expand_query_concepts(query_text):
    """Expand query terms with domain synonyms for semantic concept retrieval."""
    words = [w.lower() for w in re.split(r'\W+', query_text) if len(w) > 1 and w.lower() not in STOP_WORDS]
    expanded = set(words)
    for word in words:
        for key, syns in CONCEPT_SYNONYMS.items():
            if word == key or word in syns:
                expanded.add(key)
                expanded.update(syns)
    return expanded

def compute_tf_idf_vector(text_tokens, vocab):
    """Compute TF-IDF vector representation using word and n-gram term frequencies."""
    if not text_tokens or not vocab:
        return {}
    tf = {}
    for tok in text_tokens:
        tf[tok] = tf.get(tok, 0) + 1
    
    total = len(text_tokens)
    vec = {}
    for term, idf in vocab.items():
        if term in tf:
            vec[term] = (tf[term] / total) * idf
    return vec

def cosine_similarity(vec1, vec2):
    """Calculate Cosine Similarity between two term vectors: (A . B) / (||A|| * ||B||)"""
    if not vec1 or not vec2:
        return 0.0
    common_terms = set(vec1.keys()).intersection(set(vec2.keys()))
    if not common_terms:
        return 0.0
    dot = sum(vec1[t] * vec2[t] for t in common_terms)
    mag1 = math.sqrt(sum(v ** 2 for v in vec1.values()))
    mag2 = math.sqrt(sum(v ** 2 for v in vec2.values()))
    if mag1 == 0 or mag2 == 0:
        return 0.0
    return dot / (mag1 * mag2)

def semantic_rag_retrieval(book_text, query, ctx_chars=2200, top_k=3, book_id=None):
    """
    Real vector semantic search using sentence-transformers:
    Falls back to TF-IDF if book_id is not specified or no embeddings exist.
    Returns list of (chunk_text, char_offset, chapter_id)
    """
    if book_id:
        try:
            conn = get_db_conn()
            c = conn.cursor()
            c.execute("SELECT chunk_text, embedding, char_offset, chapter_id FROM chunk_embeddings WHERE book_id=%s", (book_id,))
            rows = c.fetchall()
            conn.close()
            
            if rows:
                model = get_embedding_model()
                query_emb = model.encode(query, normalize_embeddings=True)
                
                scored = []
                for chunk_text, emb_blob, offset, ch_id in rows:
                    chunk_emb = np.frombuffer(emb_blob, dtype=np.float32)
                    score = float(np.dot(query_emb, chunk_emb))
                    scored.append((score, chunk_text, offset, ch_id))
                
                scored.sort(key=lambda x: x[0], reverse=True)
                
                selected_chunks = []
                total_len = 0
                for score, text, offset, ch_id in scored[:top_k]:
                    if total_len + len(text) <= ctx_chars or not selected_chunks:
                        selected_chunks.append((text, offset, ch_id))
                        total_len += len(text) + 2
                
                return selected_chunks, scored[0][0] if scored else 0.0
        except Exception as e:
            print("Semantic search failed, falling back to TF-IDF:", e)

    # TF-IDF Fallback
    if not book_text:
        return [], -1.0

    paragraphs = [p.strip() for p in book_text.split('\n\n') if len(p.strip()) > 30]
    if not paragraphs:
        paragraphs = [book_text[i:i+600] for i in range(0, len(book_text), 500)]

    para_tokens = []
    doc_freq = {}
    N = len(paragraphs)

    for p in paragraphs:
        tokens = [w.lower() for w in re.split(r'\W+', p) if len(w) > 1]
        unique_tokens = set(tokens)
        para_tokens.append(tokens)
        for t in unique_tokens:
            doc_freq[t] = doc_freq.get(t, 0) + 1

    vocab_idf = {t: math.log((N + 1) / (df + 1)) + 1.0 for t, df in doc_freq.items()}

    query_expanded_tokens = list(expand_query_concepts(query))
    if not query_expanded_tokens:
        query_expanded_tokens = [w.lower() for w in re.split(r'\W+', query) if len(w) > 1]

    query_vec = compute_tf_idf_vector(query_expanded_tokens, vocab_idf)

    scored_paras = []
    for idx, (p, tokens) in enumerate(zip(paragraphs, para_tokens)):
        p_vec = compute_tf_idf_vector(tokens, vocab_idf)
        cos_sim = cosine_similarity(query_vec, p_vec)
        
        tok_set = set(tokens)
        overlap = sum(1 for qt in query_expanded_tokens if qt in tok_set)
        concept_score = overlap / (len(query_expanded_tokens) + 1)

        total_score = (cos_sim * 0.65) + (concept_score * 0.35)
        scored_paras.append((total_score, idx, p))

    scored_paras.sort(key=lambda x: x[0], reverse=True)

    if not scored_paras or scored_paras[0][0] <= 0:
        return [(book_text[:ctx_chars], 0, None)], 0.0

    selected_chunks = []
    current_length = 0
    best_score = scored_paras[0][0]

    for score, idx, para in scored_paras[:top_k]:
        if current_length + len(para) <= ctx_chars or not selected_chunks:
            # Estimate character offset of this paragraph in the full text
            offset = book_text.find(para)
            if offset == -1:
                offset = 0
            selected_chunks.append((para, offset, None))
            current_length += len(para) + 2

    return selected_chunks, best_score

def detailed_book_citation_search(book_text, query, toc=None, ctx_chars=2200, book_id=None):
    if not book_text:
        return "", {"chapter": "N/A", "page": "N/A", "exact_quote": "No book loaded"}
    
    query_clean = query.strip().lower()
    
    # 1. Check for explicit chapter/unit patterns (e.g. "Chapter 2")
    ch_match = re.search(r'(chapter|unit|lesson|ch|part|module)\s*#?\s*(\d+)', query_clean)
    target_ch_num = ch_match.group(2) if ch_match else None
    
    if target_ch_num:
        ch_patterns = [
            f"chapter {target_ch_num}", f"unit {target_ch_num}", f"lesson {target_ch_num}",
            f"chapter {target_ch_num}:", f"unit {target_ch_num}:", f"{target_ch_num}."
        ]
        for pat in ch_patterns:
            idx = book_text.lower().find(pat)
            if idx != -1:
                start = max(0, idx)
                end = min(len(book_text), start + ctx_chars)
                best_para = book_text[start:end]
                citation_meta = {
                    "chapter": f"Chapter {target_ch_num}",
                    "page": f"Page ~{max(1, start // 1800 + 1)}",
                    "exact_quote": best_para[:250].replace('\n', ' ') + "..."
                }
                return best_para, citation_meta

    # 2. Semantic RAG Search Engine (Vector Cosine + Concept Expansion)
    retrieved, score = semantic_rag_retrieval(book_text, query, ctx_chars=ctx_chars, top_k=2, book_id=book_id)
    if not retrieved:
        return "", {"chapter": "N/A", "page": "N/A", "exact_quote": "No matches found"}

    best_para, char_pos, matched_ch_id = retrieved[0]
    all_retrieved_text = "\n\n".join(r[0] for r in retrieved)

    # 3. Page tracking from [Page X] markers
    prefix_text = book_text[:max(0, char_pos + 200)]
    page_matches = re.findall(r'\[Page\s+(\d+)\]', prefix_text, re.IGNORECASE)
    page_num = f"Page {page_matches[-1]}" if page_matches else f"Section / Page ~{max(1, char_pos // 1800 + 1)}"

    # 4. TOC / Chapter Name Resolution
    chapter_name = f"Chapter {target_ch_num}" if target_ch_num else "Textbook Content"
    
    # Resolve exact chapter name from matched chapter_id
    if matched_ch_id:
        try:
            conn = get_db_conn()
            c = conn.cursor()
            c.execute("SELECT title FROM chapters WHERE id=%s", (matched_ch_id,))
            ch_row = c.fetchone()
            conn.close()
            if ch_row:
                chapter_name = ch_row[0]
        except Exception:
            pass
    elif toc and isinstance(toc, list):
        query_keywords = expand_query_concepts(query)
        for ch in toc:
            unit = ch.get("unit", "")
            topics = ch.get("topics", [])
            unit_lower = unit.lower()
            if target_ch_num and (f"chapter {target_ch_num}" in unit_lower or f"unit {target_ch_num}" in unit_lower):
                chapter_name = unit
                break
            if any(kw in unit_lower for kw in query_keywords):
                chapter_name = unit
                break
            for tp in topics:
                if any(kw in tp.lower() for kw in query_keywords):
                    chapter_name = f"{unit} - Topic: {tp}"
                    break

    citation_meta = {
        "chapter": chapter_name,
        "page": page_num,
        "exact_quote": best_para[:250].replace('\n', ' ') + ("..." if len(best_para) > 250 else "")
    }

    return all_retrieved_text, citation_meta

def excerpt_search(book_text, query, ctx_chars=1200, book_id=None):
    if not book_text:
        return ""
    if not query or len(query.strip()) < 2:
        return book_text[:ctx_chars]
    
    retrieved, _ = semantic_rag_retrieval(book_text, query, ctx_chars=ctx_chars, top_k=2, book_id=book_id)
    if isinstance(retrieved, list):
        return "\n\n".join(r[0] for r in retrieved)
    return ""



# ------------- ROUTES -------------
@app.route("/")
def home():
    return jsonify({
        "app": "Advance Education System API Backend",
        "model": GEMINI_MODEL,
        "status": "online"
    })

@app.route("/api/auth/config", methods=["GET"])
def auth_config():
    return jsonify({
        "google_client_id": os.getenv("VITE_GOOGLE_CLIENT_ID") or os.getenv("GOOGLE_CLIENT_ID") or ""
    })

# ------------- USER BOOKS & DASHBOARD ANALYTICS ENDPOINTS -------------

def save_book_record(user_id, title, file_name, raw_text, structure):
    try:
        conn = get_db_conn()
        c = conn.cursor()
        uploaded_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
        toc_json = json.dumps(structure, ensure_ascii=False)
        
        c.execute("""
            INSERT INTO books (user_id, title, file_name, raw_text, toc_json, uploaded_at)
            VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING id
        """, (user_id, title, file_name, raw_text, toc_json))
        book_id = c.fetchone()[0]

        # Insert chapters
        for idx, ch in enumerate(structure, 1):
            ch_name = ch.get("name", f"Chapter {idx}")
            ch_topics = [tp.get("name", "") for tp in ch.get("topics", []) if tp.get("name")]
            ch_text = f"{ch_name}\n" + "\n".join(ch_topics)
            c.execute("""
                INSERT INTO chapters (book_id, user_id, chapter_number, title, content_text)
                VALUES (%s, %s, %s, %s, %s) RETURNING id
            """, (book_id, user_id, str(idx), ch_name, ch_text))
            ch_id = c.fetchone()[0]

            # Initial zero-mastery record for every chapter
            c.execute("""
                INSERT INTO mastery (user_id, book_id, chapter_id, chapter_title, mastery_score, last_reviewed_at, total_attempts, correct_count)
                VALUES (%s, %s, %s, %s, 0.0, NOW(), 0, 0)
                ON CONFLICT (user_id, chapter_id) DO NOTHING
            """, (user_id, book_id, ch_id, ch_name))

        # Generate and store embeddings for semantic search
        try:
            generate_and_store_embeddings(book_id, raw_text, conn)
        except Exception as emb_e:
            print("Embedding generation failed, continuing:", emb_e)

        conn.commit()
        conn.close()
        return book_id
    except Exception as e:
        print("save_book_record notice:", e)
        return None

@app.route("/api/user/books", methods=["GET"])
def get_user_books():
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"error": "Unauthorized session"}), 401
    
    conn = get_db_conn()
    c = conn.cursor()
    c.execute("""
        SELECT b.id, b.title, b.file_name, b.uploaded_at, COUNT(ch.id) as chapter_count
        FROM books b
        LEFT JOIN chapters ch ON ch.book_id = b.id
        WHERE b.user_id = %s
        GROUP BY b.id, b.title, b.file_name, b.uploaded_at
        ORDER BY b.id DESC
    """, (user_id,))
    rows = c.fetchall()
    conn.close()

    books = [
        {
            "id": r[0],
            "title": r[1],
            "file_name": r[2],
            "uploaded_at": r[3],
            "chapter_count": r[4]
        } for r in rows
    ]
    return jsonify({"success": True, "books": books})

@app.route("/api/user/books/<int:book_id>/delete", methods=["POST", "DELETE"])
def delete_user_book(book_id):
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"error": "Unauthorized session"}), 401
    
    conn = get_db_conn()
    c = conn.cursor()
    c.execute("DELETE FROM books WHERE id=%s AND user_id=%s", (book_id, user_id))
    conn.commit()
    conn.close()
    return jsonify({"success": True, "message": "Book deleted successfully"})

@app.route("/api/user/books/active", methods=["POST"])
def set_active_book():
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"error": "Not authenticated"}), 401

    data = request.get_json() or {}
    book_id = data.get("book_id")
    if not book_id:
        return jsonify({"error": "book_id parameter required"}), 400

    conn = get_db_conn()
    c = conn.cursor()
    c.execute("SELECT id, title, raw_text, toc_json FROM books WHERE id=%s AND user_id=%s", (book_id, user_id))
    row = c.fetchone()

    if not row:
        conn.close()
        return jsonify({"error": "Book not found"}), 404

    book_text = row[2] or ""
    try:
        toc = json.loads(row[3] or "[]")
    except Exception:
        toc = []

    session["active_book_id"] = book_id
    save_memory(user_id, book_text=book_text, pastpaper_text="", toc=toc)
    conn.close()

    return jsonify({
        "success": True,
        "message": f"'{row[1]}' set as active book!",
        "active_book": {
            "id": row[0],
            "title": row[1],
            "toc": toc
        }
    })

# ------------- DEDICATED BOOK UPLOAD ENDPOINT -------------
MAX_UPLOAD_SIZE_MB = 30
ALLOWED_EXTENSIONS = {".pdf", ".docx"}

@app.route("/api/user/books/upload", methods=["POST"])
def upload_user_book():
    """Upload a PDF/DOCX file, extract text, save to books+chapters, set as active."""
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"error": "Not authenticated"}), 401

    book_file = request.files.get("book")
    if not book_file or not book_file.filename:
        return jsonify({"error": "No file uploaded. Please select a PDF or DOCX file."}), 400

    # Validate file extension
    filename = book_file.filename
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"Unsupported file type '{ext}'. Only .pdf and .docx files are allowed."}), 400

    # Validate file size (read content-length header or measure stream)
    book_file.seek(0, 2)  # seek to end
    file_size = book_file.tell()
    book_file.seek(0)  # reset to start
    max_bytes = MAX_UPLOAD_SIZE_MB * 1024 * 1024
    if file_size > max_bytes:
        size_mb = round(file_size / (1024 * 1024), 1)
        return jsonify({"error": f"File too large ({size_mb}MB). Maximum allowed size is {MAX_UPLOAD_SIZE_MB}MB."}), 413

    # Save to temp file and extract text using existing extract_text_any()
    tmp_path = os.path.join(tempfile.gettempdir(), filename)
    book_file.save(tmp_path)
    
    try:
        book_text = extract_text_any(tmp_path)
    except Exception as ext_e:
        try:
            os.remove(tmp_path)
        except OSError:
            pass
        return jsonify({"error": "This file couldn't be read — it may be corrupted or not a valid PDF/DOCX. Please try a different file."}), 400

    if not book_text or len(book_text.strip()) < 50:
        try:
            os.remove(tmp_path)
        except OSError:
            pass
        return jsonify({"error": "Could not extract meaningful text from this file. The file may be empty, scanned without OCR text, or corrupted."}), 400

    try:

        # Extract structure (TOC)
        structure = extract_structure_any(book_text) or []

        # Custom title from form or filename
        title = request.form.get("title") or os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title()

        # Save to books + chapters tables
        book_id = save_book_record(user_id, title, filename, book_text, structure)
        if not book_id:
            return jsonify({"error": "Failed to save book to database."}), 500

        # Build TOC for response and teacher memory
        toc = []
        for ch in structure:
            unit_name = ch.get("name", "")
            topics = [tp.get("name", "") or "" for tp in ch.get("topics", [])]
            toc.append({"unit": unit_name, "topics": topics})

        # Set as active book in session + teacher memory
        session["active_book_id"] = book_id
        save_memory(user_id, book_text=book_text, pastpaper_text="", toc=toc)

        # Clean up temp file
        try:
            os.remove(tmp_path)
        except OSError:
            pass

        return jsonify({
            "success": True,
            "book_id": book_id,
            "title": title,
            "file_name": filename,
            "chapter_count": len(structure),
            "toc": toc,
            "message": f"'{title}' uploaded and saved to your library!"
        })

    except Exception as e:
        return jsonify({"error": f"Upload processing failed: {str(e)}"}), 500

@app.route("/api/user/dashboard-stats", methods=["GET"])
def get_dashboard_stats():
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"error": "Not authenticated"}), 401

    conn = get_db_conn()
    c = conn.cursor()

    # 1. User streak & info
    c.execute("SELECT streak_days, last_active_date FROM users WHERE id=%s", (user_id,))
    u_row = c.fetchone()
    streak_days = u_row[0] if u_row and u_row[0] else 1

    # 2. Uploaded books
    c.execute("""
        SELECT b.id, b.title, b.file_name, b.uploaded_at, COUNT(ch.id) as chapter_count
        FROM books b
        LEFT JOIN chapters ch ON ch.book_id = b.id
        WHERE b.user_id = %s
        GROUP BY b.id, b.title, b.file_name, b.uploaded_at
        ORDER BY b.id DESC
    """, (user_id,))
    rows = c.fetchall()

    # Legacy auto-migration check: if user has memory but zero books in books table, convert memory into a saved book!
    if not rows:
        mem_book, _, mem_toc = get_memory(user_id)
        if not mem_book:
            mem_book, _, mem_toc = get_memory("student")
        if mem_book:
            structure = extract_structure_any(mem_book) or []
            auto_title = "Computer Architecture & Systems Textbook"
            new_b_id = save_book_record(user_id, auto_title, "textbook_archive.pdf", mem_book, structure)
            if new_b_id:
                c.execute("""
                    SELECT b.id, b.title, b.file_name, b.uploaded_at, COUNT(ch.id) as chapter_count
                    FROM books b
                    LEFT JOIN chapters ch ON ch.book_id = b.id
                    WHERE b.user_id = ?
                    GROUP BY b.id
                    ORDER BY b.id DESC
                """, (user_id,))
                rows = c.fetchall()

    books = [
        {"id": r[0], "title": r[1], "file_name": r[2], "uploaded_at": r[3], "chapter_count": r[4]}
        for r in rows
    ]

    # 3. Overall average mastery %
    c.execute("SELECT AVG(mastery_score) FROM mastery WHERE user_id=%s", (user_id,))
    avg_row = c.fetchone()
    overall_mastery = round(avg_row[0], 1) if (avg_row and avg_row[0] is not None) else 0.0

    # 4. Weak topics (chapters sorted ascending by mastery score)
    c.execute("""
        SELECT m.chapter_id, m.chapter_title, m.mastery_score, m.total_attempts, b.title as book_title
        FROM mastery m
        LEFT JOIN books b ON b.id = m.book_id
        WHERE m.user_id = %s
        ORDER BY m.mastery_score ASC
        LIMIT 5
    """, (user_id,))
    weak_topics = [
        {
            "chapter_id": r[0],
            "title": r[1],
            "mastery_score": round(r[2], 1),
            "total_attempts": r[3],
            "book_title": r[4] or "Textbook"
        } for r in c.fetchall()
    ]

    # 5. Recent quiz submissions
    c.execute("""
        SELECT s.id, s.score, s.total, s.percentage, s.attempted_at, b.title as book_title
        FROM quiz_submissions s
        LEFT JOIN books b ON b.id = s.book_id
        WHERE s.user_id = %s
        ORDER BY s.id DESC
        LIMIT 5
    """, (user_id,))
    recent_attempts = [
        {
            "id": r[0],
            "score": r[1],
            "total": r[2],
            "percentage": round(r[3], 1),
            "attempted_at": r[4],
            "book_title": r[5] or "Quiz"
        } for r in c.fetchall()
    ]

    # 6. Generated study materials count
    c.execute("SELECT COUNT(*) FROM study_materials WHERE user_id=%s", (user_id,))
    mat_count = c.fetchone()[0]

    conn.close()

    return jsonify({
        "authenticated": True,
        "user_id": user_id,
        "books": books,
        "overall_mastery": overall_mastery,
        "weak_topics": weak_topics,
        "recent_attempts": recent_attempts,
        "streak_days": streak_days,
        "materials_count": mat_count
    })

@app.route("/api/quiz/save-attempt", methods=["POST"])
def save_quiz_attempt():
    user_id = session.get("user_id") or (request.get_json() or {}).get("user_id")
    if is_guest_user(user_id):
        return jsonify({"success": True, "submission_id": 0, "note": "Guest quiz attempt not persisted to DB"})

    data = request.get_json() or {}
    quiz_id = data.get("quiz_id")
    book_id = data.get("book_id")
    chapter_id = data.get("chapter_id")
    chapter_title = data.get("chapter_title") or "Chapter Quiz"
    score = int(data.get("score") or 0)
    total = int(data.get("total") or 0)
    answers = data.get("answers") or []

    percentage = (score / total * 100.0) if total > 0 else 0.0
    attempted_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")

    conn = get_db_conn()
    c = conn.cursor()

    # 1. Save quiz summary submission
    c.execute("""
        INSERT INTO quiz_submissions (user_id, quiz_id, book_id, chapter_id, score, total, percentage, attempted_at)
        VALUES (%s, %s, %s, %s, %s, %s, %s, NOW()) RETURNING id
    """, (user_id, quiz_id, book_id, chapter_id, score, total, percentage))
    submission_id = c.fetchone()[0]

    # 2. Save detailed per-question attempts
    for a in answers:
        c.execute("""
            INSERT INTO attempts (user_id, quiz_id, question_id, question_text, selected_answer, correct_answer, is_correct, attempted_at)
            VALUES (%s, %s, %s, %s, %s, %s, %s, NOW())
        """, (
            user_id, quiz_id, str(a.get("id")), a.get("question"),
            str(a.get("given")), str(a.get("expected")),
            True if a.get("correct") else False
        ))

    # 3. Upsert Mastery record (UNIQUE on user_id, chapter_id)
    if chapter_id:
        c.execute("""
            INSERT INTO mastery (user_id, book_id, chapter_id, chapter_title, mastery_score, last_reviewed_at, total_attempts, correct_count)
            VALUES (%s, %s, %s, %s, %s, NOW(), 1, %s)
            ON CONFLICT (user_id, chapter_id) DO UPDATE SET
                correct_count = mastery.correct_count + EXCLUDED.correct_count,
                total_attempts = mastery.total_attempts + 1,
                mastery_score = ((mastery.correct_count + EXCLUDED.correct_count) * 100.0) / (mastery.total_attempts + 1),
                last_reviewed_at = EXCLUDED.last_reviewed_at
        """, (user_id, book_id, chapter_id, chapter_title, percentage, score))

    # 4. Update user streak
    today_str = datetime.utcnow().strftime("%Y-%m-%d")
    c.execute("SELECT last_active_date, streak_days FROM users WHERE id=%s", (user_id,))
    u_row = c.fetchone()
    if u_row:
        last_active, current_streak = u_row[0], u_row[1] or 1
        if str(last_active) != today_str:
            new_streak = current_streak + 1 if last_active else 1
            c.execute("UPDATE users SET last_active_date=NOW(), streak_days=%s WHERE id=%s", (new_streak, user_id))

    conn.commit()
    conn.close()

    return jsonify({"success": True, "submission_id": submission_id, "percentage": percentage})

@app.route("/api/user/study-materials", methods=["GET", "POST"])
def manage_study_materials():
    user_id = session.get("user_id") or (request.get_json() or {}).get("user_id")
    if is_guest_user(user_id):
        if request.method == "POST":
            return jsonify({"success": True, "id": 0, "note": "Guest material not persisted to DB"})
        return jsonify({"success": True, "materials": []})

    conn = get_db_conn()
    c = conn.cursor()

    if request.method == "POST":
        data = request.get_json() or {}
        book_id = data.get("book_id")
        mat_type = data.get("type", "study_asset")
        title = data.get("title", "Study Material")
        content_or_path = data.get("content_or_path", "")
        generated_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")

        c.execute("""
            INSERT INTO study_materials (user_id, book_id, type, title, content_or_path, generated_at)
            VALUES (%s, %s, %s, %s, %s, NOW()) RETURNING id
        """, (user_id, book_id, mat_type, title, content_or_path))
        mat_id = c.fetchone()[0]
        conn.commit()
        conn.close()
        return jsonify({"success": True, "id": mat_id})

    # GET request
    c.execute("""
        SELECT id, book_id, type, title, content_or_path, generated_at
        FROM study_materials WHERE user_id=%s ORDER BY id DESC
    """, (user_id,))
    materials = [
        {"id": r[0], "book_id": r[1], "type": r[2], "title": r[3], "content_or_path": r[4], "generated_at": r[5]}
        for r in c.fetchall()
    ]
    conn.close()
    return jsonify({"success": True, "materials": materials})

# ------------- AUTHENTICATION ENDPOINTS -------------
def verify_google_token(token):
    try:
        resp = requests.get(f"https://oauth2.googleapis.com/tokeninfo?id_token={token}", timeout=10)
        if resp.status_code == 200:
            return resp.json()
    except Exception as e:
        print("Google tokeninfo check error:", e)
    
    try:
        parts = token.split(".")
        if len(parts) == 3:
            import base64
            padded = parts[1] + "=" * (-len(parts[1]) % 4)
            decoded = base64.b64decode(padded)
            return json.loads(decoded)
    except Exception as e:
        print("JWT decode fallback error:", e)
    return None

def migrate_legacy_user_data(old_id, new_id):
    if not old_id or not new_id or old_id == new_id:
        return
    try:
        conn = get_db_conn()
        c = conn.cursor()
        c.execute("SELECT user_id FROM teacher_memory WHERE user_id=%s", (new_id,))
        if not c.fetchone():
            c.execute("UPDATE teacher_memory SET user_id=%s WHERE user_id=%s", (new_id, old_id))
        c.execute("UPDATE quizzes SET user_id=%s WHERE user_id=%s", (new_id, old_id))
        c.execute("UPDATE messages SET user_id=%s WHERE user_id=%s", (new_id, old_id))
        conn.commit()
        conn.close()
    except Exception as e:
        print("Legacy data migration notice:", e)

@app.route("/api/auth/signup", methods=["POST"])
def auth_signup():
    try:
        data = request.get_json() or {}
        username = (data.get("username") or data.get("name") or "").strip()
        email = (data.get("email") or "").strip().lower()
        password = data.get("password") or ""

        if not email or not password or not username:
            return jsonify({"error": "Full name, email address, and password are required."}), 400

        if len(password) < 4:
            return jsonify({"error": "Password must be at least 4 characters long."}), 400

        conn = get_db_conn()
        c = conn.cursor()
        c.execute("SELECT id FROM users WHERE email=%s", (email,))
        if c.fetchone():
            conn.close()
            return jsonify({"error": "An account with this email address already exists. Please log in instead."}), 400

        user_id = "usr_" + datetime.utcnow().strftime("%Y%m%d%H%M%S") + "_" + os.urandom(3).hex()
        pw_hash = generate_password_hash(password)
        created_at = datetime.utcnow().isoformat()

        c.execute("""
            INSERT INTO users (id, username, email, password_hash, picture, provider, created_at)
            VALUES (%s, %s, %s, %s, %s, %s, NOW())
        """, (user_id, username, email, pw_hash, "", "email"))
        conn.commit()
        conn.close()

        session["user_id"] = user_id
        session.permanent = True

        return jsonify({
            "success": True,
            "user": {
                "id": user_id,
                "name": username,
                "email": email,
                "picture": "",
                "provider": "email"
            }
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/auth/login", methods=["POST"])
def auth_login():
    try:
        data = request.get_json() or {}
        email = (data.get("email") or "").strip().lower()
        password = data.get("password") or ""

        if not email or not password:
            return jsonify({"error": "Email address and password are required."}), 400

        conn = get_db_conn()
        c = conn.cursor()
        c.execute("SELECT id, username, email, password_hash, picture, provider FROM users WHERE email=%s", (email,))
        row = c.fetchone()
        conn.close()

        if not row:
            return jsonify({"error": "No account found with this email. Please check your credentials or sign up."}), 401

        user_id, username, email_val, pw_hash, picture, provider = row
        if not pw_hash or not check_password_hash(pw_hash, password):
            return jsonify({"error": "Incorrect password. Please try again."}), 401

        session["user_id"] = user_id
        session.permanent = True

        return jsonify({
            "success": True,
            "user": {
                "id": user_id,
                "name": username,
                "email": email_val,
                "picture": picture or "",
                "provider": provider or "email"
            }
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/auth/google", methods=["POST"])
def auth_google():
    try:
        data = request.get_json() or {}
        credential = data.get("credential") or data.get("token") or ""
        
        payload = None
        if credential:
            payload = verify_google_token(credential)

        if not payload and "email" in data:
            payload = data

        if not payload or "email" not in payload:
            return jsonify({"error": "Invalid Google login credential"}), 400

        email = payload["email"].strip().lower()
        name = payload.get("name") or payload.get("given_name") or email.split("@")[0]
        picture = payload.get("picture") or ""
        google_sub = payload.get("sub") or ""

        conn = get_db_conn()
        c = conn.cursor()
        c.execute("SELECT id, username, email, picture, provider FROM users WHERE email=%s", (email,))
        row = c.fetchone()

        created_at = datetime.utcnow().isoformat()
        if row:
            user_id = row[0]
            c.execute("UPDATE users SET username=%s, picture=%s, provider='google' WHERE id=%s", (name, picture, user_id))
        else:
            user_id = "goog_" + (google_sub or datetime.utcnow().strftime("%Y%m%d%H%M%S"))
            c.execute("""
                INSERT INTO users (id, username, email, password_hash, picture, provider, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, NOW())
            """, (user_id, name, email, "", picture, "google"))
        
        conn.commit()
        conn.close()

        session["user_id"] = user_id
        session.permanent = True

        return jsonify({
            "success": True,
            "user": {
                "id": user_id,
                "name": name,
                "email": email,
                "picture": picture,
                "provider": "google",
                "verified": True
            }
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/auth/me", methods=["GET"])
def auth_me():
    user_id = session.get("user_id")
    if not user_id:
        return jsonify({"authenticated": False}), 401
    
    conn = get_db_conn()
    c = conn.cursor()
    c.execute("SELECT id, username, email, picture, provider FROM users WHERE id=%s", (user_id,))
    row = c.fetchone()
    conn.close()

    if not row:
        session.clear()
        return jsonify({"authenticated": False}), 401

    return jsonify({
        "authenticated": True,
        "user": {
            "id": row[0],
            "name": row[1],
            "email": row[2],
            "picture": row[3] or "",
            "provider": row[4] or "email"
        }
    })

@app.route("/api/auth/logout", methods=["POST"])
def auth_logout():
    session.clear()
    return jsonify({"success": True, "message": "Logged out successfully"})



@app.route("/train_teacher", methods=["POST"])
def train_teacher():
    try:
        user_id = session.get("user_id")
        if not user_id:
            return jsonify({"error": "Not authenticated"}), 401
        book_file = request.files.get("teacher_book")
        past_file = request.files.get("past_paper")
        if not book_file:
            return jsonify({"error":"Please upload a book file"}), 400
        tmp = os.path.join(tempfile.gettempdir(), book_file.filename)
        book_file.save(tmp)
        book_text = extract_text_any(tmp)
        past_text = ""
        if past_file:
            tmp2 = os.path.join(tempfile.gettempdir(), past_file.filename)
            past_file.save(tmp2)
            past_text = extract_text_any(tmp2)
        # extract structure (TOC)
        structure = extract_structure_any(book_text) or []
        toc = []
        for ch in structure:
            unit_name = ch.get("name","")
            topics = []
            for tp in ch.get("topics",[]):
                tname = tp.get("name","") or ""
                topics.append(tname)
            toc.append({"unit": unit_name, "topics": topics})
        
        save_memory(user_id, book_text=book_text, pastpaper_text=past_text, toc=toc)

        # Reuse active book if already uploaded, to avoid duplicate entries
        book_id = session.get("active_book_id")
        if not book_id:
            book_id = save_book_record(user_id, book_file.filename, book_file.filename, book_text, structure)
            if book_id:
                session["active_book_id"] = book_id

        # clear messages for fresh session
        conn = get_db_conn(); c = conn.cursor(); c.execute("DELETE FROM messages WHERE user_id=%s", (user_id,)); conn.commit(); conn.close()
        return jsonify({"success": True, "book_id": book_id, "toc": toc, "message": "Book trained and saved to your personal library!"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/get_topics", methods=["POST"])
def get_topics():
    try:
        data = request.get_json() or {}
        user_id = data.get("user_id", "student")
        unit_name = (data.get("unit") or "").strip()
        book, past, toc = get_memory(user_id)
        if not book:
            return jsonify({"error":"No trained book"}), 400
        for ch in toc:
            if ch.get("unit","").strip().lower() == unit_name.lower() or unit_name.lower() in ch.get("unit","").lower():
                return jsonify({"topics": ch.get("topics", [])})
        if toc:
            return jsonify({"topics": toc[0].get("topics", [])})
        return jsonify({"topics": []})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/teach_topic", methods=["POST"])
def teach_topic():
    try:
        data = request.get_json() or {}
        user_id = data.get("user_id", "student")
        topic = (data.get("topic") or "").strip()
        use_web_search = data.get("use_web_search", True)
        if not topic:
            return jsonify({"error":"Topic required"}), 400
        book_text, past_text, _ = get_memory(user_id)
        add_message(user_id, "user", f"Teach topic: {topic}")
        
        system_prompt = ("You are a patient, friendly teacher powered by Gemini AI. Teach the topic clearly. "
                         "Provide: 1) Clear explanation, 2) Two short examples, 3) One practice task.\n"
                         "Include a section '### 🧠 LLM Knowledge & Concept References' at the bottom.")
        messages = [{"role":"system","content":system_prompt}]
        
        if book_text:
            excerpt = excerpt_search(book_text, topic, ctx_chars=1600)
            messages.append({"role":"system","content":"Book excerpt:\n\n" + excerpt})
            
        web_sources = []
        if use_web_search:
            web_context, web_sources = perform_web_search(topic, max_results=4)
            if web_context:
                messages.append({"role":"system","content":web_context})

        history = get_history(user_id, limit=10)
        messages.extend(history)
        messages.append({"role":"user","content":f"Please teach: {topic}"})
        reply = call_gemini(messages, max_tokens=350, temperature=0.25)

        add_message(user_id, "assistant", reply)
        return jsonify({"reply": reply, "web_sources": web_sources})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/ai_chat", methods=["POST"])
def ai_chat():
    try:
        body = request.get_json() or {}
        user_id = body.get("user_id", "student")
        message = body.get("message","").strip()
        use_web_search = body.get("use_web_search", True)
        if not message:
            return jsonify({"error":"Message empty"}), 400
        add_message(user_id, "user", message)
        
        system_prompt = (
            "You are an expert, patient AI teacher powered by Gemini AI. "
            "Answer the user's question clearly and concisely with key points.\n"
            "At the end of your response, ALWAYS include a dedicated section titled:\n"
            "### 🧠 LLM Knowledge & Concept References\n"
            "Listing the primary theoretical principles and model reasoning used."
        )
        messages = [{"role":"system","content":system_prompt}]

        web_sources = []
        if use_web_search:
            web_context, web_sources = perform_web_search(message, max_results=3)
            if web_context:
                messages.append({"role":"system","content":web_context})

        messages.extend(get_history(user_id, limit=8))
        messages.append({"role":"user","content":message})
        reply = call_gemini(messages, max_tokens=350, temperature=0.2)
        add_message(user_id, "assistant", reply)

        return jsonify({
            "reply": reply,
            "web_sources": web_sources,
            "llm_reference": "Gemini AI Parametric Weights & Internal Reasoning"
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/clear_history", methods=["POST"])
def clear_history():
    try:
        body = request.get_json() or {}
        user_id = body.get("user_id", "student")
        if not is_guest_user(user_id):
            conn = get_db_conn(); c = conn.cursor()
            c.execute("DELETE FROM messages WHERE user_id=%s", (user_id,))
            conn.commit(); conn.close()
        return jsonify({"success": True, "message": "Chat history cleared successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/ai_chat_stream", methods=["POST"])
def ai_chat_stream():
    try:
        body = request.get_json() or {}
        user_id = body.get("user_id", "student")
        message = body.get("message","").strip()
        use_web_search = body.get("use_web_search", True)
        if not message:
            return jsonify({"error":"Message empty"}), 400

        # Fetch past history BEFORE adding current message to prevent double-logging in LLM context
        history = get_history(user_id, limit=8)
        add_message(user_id, "user", message)

        system_prompt = (
            "You are an expert, patient AI teacher powered by Gemini AI. "
            "Answer the user's question clearly and concisely with key points.\n"
            "At the end of your response, ALWAYS include a dedicated section titled:\n"
            "### 🧠 LLM Knowledge & Concept References\n"
            "Listing the primary theoretical principles and model reasoning used."
        )
        messages = [{"role":"system","content":system_prompt}]

        web_sources = []
        user_content = message
        if use_web_search:
            web_context, web_sources = perform_web_search(message, max_results=3)
            if web_context:
                user_content += f"\n\n{web_context}"

        messages.extend(history)
        messages.append({"role":"user","content":user_content})

        def generate():
            yield f"data: {json.dumps({'type':'sources', 'web_sources': web_sources, 'llm_reference':'Gemini AI Parametric Weights & Internal Reasoning'})}\n\n"
            full_text = ""
            try:
                for text_chunk in stream_gemini(messages, max_tokens=1000, temperature=0.2):
                    if text_chunk:
                        full_text += text_chunk
                        yield f"data: {json.dumps({'type':'text', 'content': text_chunk})}\n\n"
                if full_text:
                    add_message(user_id, "assistant", full_text)
            except Exception as ex:
                yield f"data: {json.dumps({'type':'error', 'error': str(ex)})}\n\n"
            yield "data: [DONE]\n\n"

        return Response(generate(), mimetype='text/event-stream')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/ask_book_teacher_stream", methods=["POST"])
def ask_book_teacher_stream():
    try:
        body = request.get_json() or {}
        user_id = body.get("user_id", "student")
        question = body.get("question", "").strip()
        mode = body.get("mode", "normal")
        panic_time = body.get("panic_time", "24 Hours")

        if not question:
            return jsonify({"error": "Question empty"}), 400

        book_text, past_text, toc = get_memory(user_id)
        if not book_text:
            return jsonify({"error": "No book uploaded yet. Please upload or paste a book first under 'Upload Book' section."}), 400

        book_id = session.get("active_book_id")
        excerpt, citation = detailed_book_citation_search(book_text, question, toc=toc, book_id=book_id)

        mode_instructions = ""
        if mode == "panic":
            mode_instructions = (
                f"CRITICAL MODE: 🚨 EMERGENCY EXAM PANIC MODE ({panic_time} left!).\n"
                f"- The student has an urgent exam deadline ({panic_time}). Give ultra-concise, high-yield bullet points.\n"
                f"- Highlight exact exam definitions, core formulas, and high-yield scoring points. ZERO fluff or filler!"
            )
        elif mode == "curiosity":
            mode_instructions = (
                "CRITICAL MODE: 🔍 CURIOSITY DEEP DIVE MODE.\n"
                "- The student wants deep conceptual mastery. Provide advanced insights, real-world industry applications, edge-cases, and deep theoretical connections beyond standard textbook definitions."
            )
        elif mode == "examprep":
            mode_instructions = (
                "CRITICAL MODE: 🎯 EXAM PREP & PAST PAPER MODE.\n"
                "- Format answers like marking schemes. Provide key exam phrases required for full marks, common student traps/mistakes to avoid, and scoring strategies."
            )
        elif mode == "feynman":
            mode_instructions = (
                "CRITICAL MODE: 🧠 FEYNMAN / ELI5 MODE.\n"
                "- Explain complex concepts using intuitive 5th-grade analogies, plain English, zero jargon, and clear mental models."
            )
        else:
            mode_instructions = "CRITICAL MODE: 📘 NORMAL STUDY MODE. Provide balanced, comprehensive, step-by-step guidance."

        system_prompt = (
            "You are an encouraging, expert Book Extractor Teacher.\n"
            f"{mode_instructions}\n\n"
            "Teach the student step-by-step using the provided Book Excerpt below.\n"
            "1. Give a clear explanation of the main concepts, definitions, and key points in the excerpt.\n"
            "2. Organize key takeaways into bullet points or numbered steps.\n"
            "Keep your response structured and focused."
        )

        user_prompt = f"Student Question: {question}\n\nStrict Book Excerpt:\n{excerpt}"

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        def generate():
            yield f"data: {json.dumps({'type':'citation', 'citation': citation})}\n\n"
            full_text = ""
            try:
                for text_chunk in stream_gemini(messages, max_tokens=1000, temperature=0.1):
                    if text_chunk:
                        full_text += text_chunk
                        yield f"data: {json.dumps({'type':'text', 'content': text_chunk})}\n\n"
                add_message(user_id, "assistant", full_text)
            except Exception as ex:
                yield f"data: {json.dumps({'type':'error', 'error': str(ex)})}\n\n"
            yield "data: [DONE]\n\n"

        return Response(generate(), mimetype='text/event-stream')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate_flashcards", methods=["POST"])
def generate_flashcards():
    try:
        body = request.get_json() or {}
        topic = body.get("topic", "").strip()
        count = int(body.get("count", 5))
        user_id = body.get("user_id", "student")
        source = body.get("source", "web_llm")

        if not topic:
            return jsonify({"error": "Topic required"}), 400

        book_text, _, _ = get_memory(user_id)
        matching_excerpt = excerpt_search(book_text, topic, ctx_chars=2000) if (book_text and source == "book") else ""

        if matching_excerpt:
            user_prompt = f"Topic: {topic}\n\nReference Material Excerpt from Book:\n{matching_excerpt}"
        else:
            user_prompt = f"Topic: {topic}"

        system_prompt = (
            f"You are an expert flashcard generator powered by Gemini AI. "
            f"Generate exactly {count} educational flashcard Q&A pairs for the topic: '{topic}'.\n"
            f"Output ONLY a valid JSON array of objects, where each object has key 'question' and key 'answer'.\n"
            f"Do not include any explanation outside the JSON array.\n"
            f"Example format: [{{\"question\": \"What is X?\", \"answer\": \"X is...\"}}]"
        )
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        raw_llm = call_gemini(messages, max_tokens=1000, temperature=0.2, response_json=True)

        
        json_match = re.search(r'\[.*\]', raw_llm, re.DOTALL)
        if json_match:
            try:
                cards = json.loads(json_match.group(0))
                formatted_cards = []
                for i, c in enumerate(cards, 1):
                    q = c.get("question") or c.get("q") or f"Question #{i} on {topic}"
                    a = c.get("answer") or c.get("a") or f"Answer #{i}"
                    formatted_cards.append({"id": i, "question": q, "answer": a, "mastered": False})
                return jsonify({"success": True, "cards": formatted_cards})
            except Exception:
                pass

        fallback_cards = [
            {
                "id": i + 1,
                "question": f"Key Question #{i + 1} regarding {topic}?",
                "answer": f"Core explanation & key concept summary for {topic} (Point #{i + 1}).",
                "mastered": False
            } for i in range(count)
        ]
        return jsonify({"success": True, "cards": fallback_cards})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate_quiz", methods=["POST"])
def generate_quiz():
    try:
        data = request.get_json() or {}
        user_id = data.get("user_id", "student")
        topic = (data.get("topic") or "").strip()
        count = int(data.get("question_count") or 5)
        quiz_type = data.get("quiz_type", "mcq")
        source = data.get("source", "web_llm")  # 'web_llm' (General Knowledge + Live Web) or 'book'
        focus_weak_areas = data.get("focus_weak_areas", False)
        
        auto_selected_topic = ""
        auto_selected_reason = ""

        if focus_weak_areas:
            book_id = session.get("active_book_id")
            if not book_id:
                # Fallback to latest uploaded book for this user
                conn = get_db_conn()
                c = conn.cursor()
                c.execute("SELECT id FROM books WHERE user_id=%s ORDER BY uploaded_at DESC LIMIT 1", (user_id,))
                row = c.fetchone()
                conn.close()
                if row:
                    book_id = row[0]
                    session["active_book_id"] = book_id
            
            if not book_id:
                return jsonify({"error": "Please upload or select a textbook first to use Weakness-First Quiz mode."}), 400
                
            conn = get_db_conn()
            c = conn.cursor()
            c.execute("SELECT id, title FROM chapters WHERE book_id=%s", (book_id,))
            all_chapters = c.fetchall()
            if not all_chapters:
                conn.close()
                return jsonify({"error": "No chapters found in the active textbook."}), 400
            
            c.execute("SELECT chapter_id, mastery_score, total_attempts, last_reviewed_at FROM mastery WHERE user_id=%s AND book_id=%s", (user_id, book_id))
            mastery_rows = {row[0]: (row[1], row[2], row[3]) for row in c.fetchall()}
            conn.close()

            # Filter attempted chapters
            attempted_chapters = []
            for chap_id, chap_title in all_chapters:
                if chap_id in mastery_rows:
                    score, attempts, reviewed_at = mastery_rows[chap_id]
                    if attempts > 0:
                        attempted_chapters.append((chap_id, chap_title, score, reviewed_at or ""))
            
            if attempted_chapters:
                # Sort attempted chapters:
                # Primary key: score ascending (lowest first)
                # Secondary key: reviewed_at ascending (oldest first). Empty string reviews naturally sort first.
                attempted_chapters.sort(key=lambda x: (x[2], x[3]))
                selected_chap_id, selected_chap_title, selected_score, _ = attempted_chapters[0]
                auto_selected_topic = selected_chap_title
                mastery_percentage = int(selected_score * 100) if selected_score is not None else 0
                auto_selected_reason = f"your current mastery here is {mastery_percentage}%"
            else:
                # Fallback: pick the first chapter in the book
                selected_chap_id, selected_chap_title = all_chapters[0]
                auto_selected_topic = selected_chap_title
                auto_selected_reason = "no attempts made yet on this book"
            
            topic = auto_selected_topic
            source = "book"

        if not topic:
            return jsonify({"error": "Topic parameter required."}), 400

        user_prompt_content = f"Topic requested: {topic}"

        # If user explicitly selected 'book', search uploaded textbook memory
        if source == "book":
            book_text, past_text, _ = get_memory(user_id)
            matching_excerpt = excerpt_search(book_text, topic, ctx_chars=2200) if book_text else ""
            if matching_excerpt:
                user_prompt_content += f"\n\nReference Material Excerpt from Uploaded Book:\n{matching_excerpt}"

        # Perform DuckDuckGo live web search to empower LLM with up-to-date facts
        web_context, web_sources = perform_web_search(topic, max_results=3)
        if web_context:
            user_prompt_content += f"\n\nReal-Time Web Search Knowledge:\n{web_context}"

        mode = data.get("mode", "normal")
        panic_time = data.get("panic_time", "24 Hours")

        mode_instructions = ""
        if mode == "panic":
            mode_instructions = f"MODE: 🚨 EMERGENCY EXAM CRUNCH ({panic_time} left). Focus on high-yield, rapid-fire exam questions testing key formulas, definitions & quick scoring points!"
        elif mode == "curiosity":
            mode_instructions = "MODE: 🔍 CURIOSITY DEEP DIVE. Focus on scenario-based, conceptual questions testing deep connections & advanced real-world applications!"
        elif mode == "examprep":
            mode_instructions = "MODE: 🎯 EXAM PREP & PAST PAPERS. Format questions like official past papers with marking scheme style options and tricky distractor choices."
        elif mode == "feynman":
            mode_instructions = "MODE: 🧠 FEYNMAN ELI5. Focus on testing intuitive conceptual understanding using clear everyday analogies."

        if quiz_type == "mcq":
            json_schema = (
                f"{{\n"
                f"  \"title\": \"Quiz on {topic}\",\n"
                f"  \"questions\": [\n"
                f"    {{\n"
                f"      \"id\": 1,\n"
                f"      \"type\": \"mcq\",\n"
                f"      \"question\": \"Targeted question about {topic}?\",\n"
                f"      \"options\": [\"A) Choice 1\", \"B) Choice 2\", \"C) Choice 3\", \"D) Choice 4\"],\n"
                f"      \"answer\": \"A\",\n"
                f"      \"explain\": \"Detailed scientific/academic explanation why option A is correct.\"\n"
                f"    }}\n"
                f"  ]\n"
                f"}}\n"
            )
        else:
            json_schema = (
                f"{{\n"
                f"  \"title\": \"Quiz on {topic}\",\n"
                f"  \"questions\": [\n"
                f"    {{\n"
                f"      \"id\": 1,\n"
                f"      \"type\": \"{quiz_type}\",\n"
                f"      \"question\": \"Conceptual or open-ended question about {topic} requiring a short answer?\",\n"
                f"      \"answer\": \"The expected correct conceptual answer/key points to look for.\",\n"
                f"      \"explain\": \"Detailed explanation of the correct concept.\"\n"
                f"    }}\n"
                f"  ]\n"
                f"}}\n"
                f"CRITICAL: Do NOT generate an 'options' field for {quiz_type} questions."
            )

        system_prompt = (
            f"You are an expert quiz generator powered by LLM and live web intelligence.\n"
            f"{mode_instructions}\n"
            f"Generate a high-quality {count}-question quiz STRICTLY focused on the topic: '{topic}'.\n"
            f"CRITICAL RULES:\n"
            f"1. Every single question MUST be directly about '{topic}'. Do NOT generate questions on unrelated subjects.\n"
            f"2. Use full LLM general knowledge and live web context to create accurate, engaging, educational questions.\n"
            f"Output ONLY valid JSON matching this exact structure:\n"
            f"{json_schema}\n"
            f"Generate exactly {count} questions."
        )


        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt_content}
        ]

        raw = call_gemini(messages, max_tokens=1500, temperature=0.3, response_json=True)
        clean_raw = re.sub(r'```(?:json)?', '', raw).strip()
        start = clean_raw.find("{")
        end = clean_raw.rfind("}")
        
        if start != -1 and end != -1:
            try:
                jtxt = clean_raw[start:end+1]
                quiz_obj = json.loads(jtxt)
                save_quiz(user_id, quiz_obj)
                res_data = {"quiz": quiz_obj}
                if focus_weak_areas:
                    res_data["auto_selected_topic"] = auto_selected_topic
                    res_data["auto_selected_reason"] = auto_selected_reason
                return jsonify(res_data)
            except Exception as pe:
                print("Quiz JSON parse warning:", pe)


        # Topic-focused fallback if JSON generation has any syntax issues
        fallback_questions = []
        for i in range(1, count + 1):
            if quiz_type == "mcq":
                fallback_questions.append({
                    "id": i,
                    "type": quiz_type,
                    "question": f"Question {i}: What is a core principle of {topic}?",
                    "options": [
                        f"A) Primary key mechanism defining {topic}",
                        f"B) Secondary process related to {topic}",
                        f"C) Incorrect hypothesis about {topic}",
                        f"D) Distractor option"
                    ],
                    "answer": "A",
                    "explain": f"Option A accurately highlights key principles of {topic}."
                })
            else:
                fallback_questions.append({
                    "id": i,
                    "type": quiz_type,
                    "question": f"Question {i}: Describe the core principle of {topic} and how it functions.",
                    "answer": f"The core principle of {topic} revolves around its fundamental mechanisms.",
                    "explain": f"This question asks you to explain the key elements of {topic}."
                })
            
        quiz_obj = {
            "title": f"Exam Quiz: {topic}",
            "questions": fallback_questions
        }
        save_quiz(user_id, quiz_obj)
        res_data = {"quiz": quiz_obj}
        if focus_weak_areas:
            res_data["auto_selected_topic"] = auto_selected_topic
            res_data["auto_selected_reason"] = auto_selected_reason
        return jsonify(res_data)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def evaluate_short_answer_with_llm(question, expected, given):
    if not given.strip():
        return False
    
    # Use Gemini API to grade the open-ended response
    messages = [
        {
            "role": "system", 
            "content": (
                "You are an automated academic grading assistant. "
                "Compare the student's answer to the expected correct answer. "
                "Determine if the student has demonstrated a correct conceptual understanding of the answer. "
                "Minor wording differences or typos should be forgiven. "
                "Output ONLY 'YES' if they are correct, or 'NO' if they are incorrect. Do NOT write any other words."
            )
        },
        {
            "role": "user", 
            "content": f"Question: {question}\nExpected Answer: {expected}\nStudent's Answer: {given}"
        }
    ]
    try:
        response = call_gemini(messages, max_tokens=10, temperature=0.1).strip().upper()
        if "YES" in response:
            return True
    except Exception as e:
        print("LLM short answer evaluation error:", e)
    
    # Fallback to simple matching if LLM fails
    given_lower = given.lower()
    expected_lower = expected.lower()
    return given_lower in expected_lower or expected_lower in given_lower


@app.route("/submit_quiz", methods=["POST"])
def submit_quiz():
    try:
        body = request.get_json() or {}
        user_id = body.get("user_id","student")
        answers = body.get("answers", {})
        quiz = get_latest_quiz(user_id)
        if not quiz:
            return jsonify({"error":"No quiz available"}), 400
        questions = quiz.get("questions", [])
        total = len(questions); correct = 0; details = []; wrong = []
        for q in questions:
            qid = str(q.get("id"))
            expected = str(q.get("answer","")).strip()
            given = str(answers.get(qid,"")).strip()
            is_correct = False
            if q.get("type") == "mcq":
                if given and expected and given[0].lower() == expected[0].lower(): is_correct = True
                elif given.lower() == expected.lower(): is_correct = True
            else:
                # Open-ended / Short answer / Socratic questions evaluation
                is_correct = evaluate_short_answer_with_llm(q.get("question"), expected, given)
            if is_correct: correct += 1
            else:
                wrong.append({"id": qid, "question": q.get("question"), "expected": expected, "given": given, "explain": q.get("explain","")})
            details.append({"id": qid, "question": q.get("question"), "expected": expected, "given": given, "correct": is_correct, "explain": q.get("explain","")})
        percent = (correct/total*100) if total else 0
        reteach_text = ""
        if percent < 70 and wrong:
            book_text, _, _ = get_memory(user_id)
            miss_summary = "\n".join([f"Q{w['id']}: {w['question']} (expected: {w['expected']}, given: {w['given']})" for w in wrong])
            system_prompt = ("You are a patient teacher. Student missed these. Provide a short reteach: core concept, 2 examples, 2 tiny tasks.")
            messages = [{"role":"system","content":system_prompt},
                        {"role":"system","content":"Material (truncated):\n" + (book_text[:4000] if book_text else "General knowledge")},
                        {"role":"user","content":"Missed:\n" + miss_summary}]
            reteach_text = call_gemini(messages, max_tokens=600, temperature=0.2)
        return jsonify({"score": correct, "total": total, "percentage": percent, "details": details, "reteach": reteach_text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generate_curriculum", methods=["POST"])
def generate_curriculum():
    try:
        user_id = session.get("user_id")
        if not user_id:
            return jsonify({"error": "Not authenticated"}), 401
        subject = request.form.get("subject", "General Subject")
        grade = request.form.get("grade", "9")
        language = request.form.get("language", "English")
        start_date_str = request.form.get("start_date")
        duration_val = request.form.get("duration_val", 9)
        duration_unit = request.form.get("duration_unit", "months")
        
        book_file = request.files.get("book") or request.files.get("teacher_book")
        if not book_file:
            return jsonify({"error": "Please upload a textbook file (PDF/DOCX/TXT) or paste book text."}), 400

        if not start_date_str:
            return jsonify({"error": "Start date required (YYYY-MM-DD)."}), 400

        start_date = datetime.strptime(start_date_str, "%Y-%m-%d")

        tmp = os.path.join(tempfile.gettempdir(), book_file.filename or "uploaded_book.txt")
        book_file.save(tmp)
        book_text = extract_text_any(tmp)

        structure = extract_structure_any(book_text) or []
        toc = []
        for ch in structure:
            unit_name = ch.get("name", "")
            topics = []
            for tp in ch.get("topics", []):
                tname = tp.get("name", "") or ""
                topics.append(tname)
            toc.append({"unit": unit_name, "topics": topics})

        save_memory(user_id, book_text=book_text, pastpaper_text="", toc=toc)

        # Reuse active book if already uploaded, to avoid duplicate entries
        book_id = session.get("active_book_id")
        if not book_id:
            book_id = save_book_record(user_id, f"{subject} Textbook", book_file.filename or "uploaded_book.txt", book_text, structure)
            if book_id:
                session["active_book_id"] = book_id

        df = generate_curriculum_custom(start_date, structure, subject, grade, duration_val=duration_val, duration_unit=duration_unit)
        base = f"{subject.replace(' ','_')}_Grade{grade}_{start_date.strftime('%Y%m%d')}_{duration_val}{duration_unit[:1].upper()}"
        excel_path = os.path.join(OUTPUT_FOLDER, base + ".xlsx")
        word_path = os.path.join(OUTPUT_FOLDER, base + ".docx")
        df.to_excel(excel_path, index=False)
        highlight_excel(excel_path, df)
        export_to_word(df, word_path)

        return jsonify({
            "success": True,
            "book_id": book_id,
            "excel": f"/download/{os.path.basename(excel_path)}",
            "word": f"/download/{os.path.basename(word_path)}",
            "toc": toc
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate_topic_curriculum", methods=["POST"])
def generate_topic_curriculum():

    try:
        data = request.get_json() or {}
        user_id = data.get("user_id","student")
        topic = (data.get("topic") or "").strip()
        subject = data.get("subject","Subject")
        grade = data.get("grade","Grade")
        start_date_str = data.get("start_date")
        if not start_date_str:
            return jsonify({"error":"start_date required YYYY-MM-DD"}), 400
        start_date = datetime.strptime(start_date_str, "%Y-%m-%d")
        structure = [{'name': topic or 'Topic', 'topics':[{'name': topic or 'Topic', 'subtopics':[topic or 'Topic'] }]}]
        df = generate_curriculum_9_months(start_date, structure, subject, grade)
        base = f"{subject}_Grade{grade}_{start_date.strftime('%Y%m%d')}_{(topic or 'topic').replace(' ','_')}_9M"
        excel_path = os.path.join(OUTPUT_FOLDER, base + ".xlsx")
        word_path = os.path.join(OUTPUT_FOLDER, base + ".docx")
        df.to_excel(excel_path, index=False)
        highlight_excel(excel_path, df)
        export_to_word(df, word_path)
        return jsonify({"excel": f"/download/{os.path.basename(excel_path)}", "word": f"/download/{os.path.basename(word_path)}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

import base64

def save_base64_image(base64_str):
    """Save base64 image data URL to a temporary PNG file for PPT insertion."""
    try:
        if "," in base64_str:
            base64_str = base64_str.split(",")[1]
        img_data = base64.b64decode(base64_str)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_data)
        tmp.close()
        return tmp.name
    except Exception as e:
        print("Base64 save error:", e)
        return None

def create_concept_banner_image(title_text, style_name):
    """Create a visual concept banner image for presentation slides."""
    try:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (800, 500), color=(15, 23, 42))
        draw = ImageDraw.Draw(img)
        for i in range(500):
            r = int(15 + (i / 500) * 25)
            g = int(23 + (i / 500) * 35)
            b = int(42 + (i / 500) * 65)
            draw.line([(0, i), (800, i)], fill=(r, g, b))
        
        draw.rounded_rectangle([30, 30, 770, 470], radius=18, outline=(99, 102, 241), width=3)
        draw.ellipse([580, -40, 820, 200], fill=(56, 189, 248), outline=(56, 189, 248))
        draw.ellipse([-40, 280, 260, 580], fill=(168, 85, 247), outline=(168, 85, 247))
        
        tmp_img = os.path.join(tempfile.gettempdir(), f"ppt_banner_{abs(hash(title_text)) % 100000}.png")
        img.save(tmp_img)
        return tmp_img
    except Exception as e:
        print("Banner creation warning:", e)
        return None

def build_pptx_presentation(topic, slides, color_scheme='indigo', template_style='modern_cards', include_images=True, image_style='minimalist', slide_shape='rounded_card', slide_animation='fade_glow'):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    PALETTES = {
        'indigo': {
            'bg': RGBColor(15, 23, 42),
            'card': RGBColor(30, 41, 59),
            'header': RGBColor(167, 139, 250),
            'text': RGBColor(248, 250, 252),
            'accent': RGBColor(56, 189, 248),
            'sub': RGBColor(203, 213, 225)
        },
        'cyan': {
            'bg': RGBColor(8, 51, 68),
            'card': RGBColor(21, 94, 117),
            'header': RGBColor(6, 182, 212),
            'text': RGBColor(240, 253, 250),
            'accent': RGBColor(52, 211, 153),
            'sub': RGBColor(165, 243, 252)
        },
        'emerald': {
            'bg': RGBColor(2, 44, 34),
            'card': RGBColor(6, 78, 59),
            'header': RGBColor(16, 185, 129),
            'text': RGBColor(236, 253, 245),
            'accent': RGBColor(251, 191, 36),
            'sub': RGBColor(167, 243, 208)
        },
        'amber': {
            'bg': RGBColor(69, 26, 3),
            'card': RGBColor(120, 53, 15),
            'header': RGBColor(245, 158, 11),
            'text': RGBColor(255, 251, 235),
            'accent': RGBColor(244, 63, 94),
            'sub': RGBColor(253, 230, 138)
        },
        'corporate_white': {
            'bg': RGBColor(248, 250, 252),
            'card': RGBColor(255, 255, 255),
            'header': RGBColor(15, 23, 42),
            'text': RGBColor(30, 41, 59),
            'accent': RGBColor(37, 99, 235),
            'sub': RGBColor(71, 85, 105)
        }
    }
    
    # Map slide_shape geometry
    SHAPE_MAP = {
        'circle': MSO_SHAPE.OVAL,
        'hexagon': MSO_SHAPE.HEXAGON,
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_card': MSO_SHAPE.ROUNDED_RECTANGLE
    }
    chosen_mso_shape = SHAPE_MAP.get(slide_shape, MSO_SHAPE.ROUNDED_RECTANGLE)

    colors = PALETTES.get(color_scheme, PALETTES['indigo'])
    blank_layout = prs.slide_layouts[6]
    
    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg_shape = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['bg']
    bg_shape.line.fill.background()
    
    card_shape = title_slide.shapes.add_shape(chosen_mso_shape, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = colors['card']
    card_shape.line.color.rgb = colors['accent']

    tf = card_shape.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "ADVANCE EDUCATION SYSTEM PRESENTATION"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = colors['accent']
    p0.alignment = PP_ALIGN.CENTER
    
    p1 = tf.add_paragraph()
    p1.text = topic
    p1.font.size = Pt(34)
    p1.font.bold = True
    p1.font.color.rgb = colors['text']
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"Template: {template_style.replace('_',' ').title()}  •  Palette: {color_scheme.replace('_',' ').title()}  •  Animation: {slide_animation.replace('_',' ').title()}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = colors['sub']
    p2.alignment = PP_ALIGN.CENTER

    # Content Slides
    for idx, slide_data in enumerate(slides, 1):
        s = prs.slides.add_slide(blank_layout)
        
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors['bg']
        bg.line.fill.background()
        
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = colors['card']
        hdr.line.color.rgb = colors['header']
        
        htf = hdr.text_frame
        htf.word_wrap = True
        hp0 = htf.paragraphs[0]
        hp0.text = f"SLIDE {idx} OF {len(slides)}  •  {slide_data.get('title', 'Slide Title')}"
        hp0.font.size = Pt(20)
        hp0.font.bold = True
        hp0.font.color.rgb = colors['header']
        
        if slide_data.get('subtitle'):
            hp1 = htf.add_paragraph()
            hp1.text = slide_data.get('subtitle')
            hp1.font.size = Pt(13)
            hp1.font.color.rgb = colors['sub']

        custom_img_b64 = slide_data.get('custom_image')
        has_custom_image = bool(custom_img_b64)
        has_any_image = has_custom_image or include_images
        
        content_width = Inches(7.2) if has_any_image else Inches(11.733)
        
        # Per slide chosen shape or global shape
        per_slide_shape_key = slide_data.get('shape', slide_shape)
        per_slide_mso = SHAPE_MAP.get(per_slide_shape_key, chosen_mso_shape)

        c_card = s.shapes.add_shape(per_slide_mso, Inches(0.8), Inches(1.9), content_width, Inches(5.1))
        c_card.fill.solid()
        c_card.fill.fore_color.rgb = colors['card']
        c_card.line.color.rgb = colors['accent']
        
        ctf = c_card.text_frame
        ctf.word_wrap = True
        
        bullets = slide_data.get('bullets', [])
        for b_idx, b_text in enumerate(bullets):
            p = ctf.paragraphs[0] if b_idx == 0 else ctf.add_paragraph()
            p.text = f"•  {b_text}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(12)
            
        # If user uploaded a custom image for this slide, insert it!
        if has_custom_image:
            custom_img_path = save_base64_image(custom_img_b64)
            if custom_img_path and os.path.exists(custom_img_path):
                s.shapes.add_picture(custom_img_path, Inches(8.2), Inches(1.9), Inches(4.333), Inches(5.1))
        elif include_images:
            img_path = create_concept_banner_image(slide_data.get('title', topic), image_style)
            if img_path and os.path.exists(img_path):
                s.shapes.add_picture(img_path, Inches(8.2), Inches(1.9), Inches(4.333), Inches(5.1))

    base_name = f"{re.sub(r'\\W+', '_', topic)[:28]}_Presentation"
    out_filename = f"{base_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    out_path = os.path.join(OUTPUT_FOLDER, out_filename)
    prs.save(out_path)
    return out_filename


@app.route("/generate_ppt_slides", methods=["POST"])
def generate_ppt_slides():
    try:
        body = request.get_json() or {}
        topic = body.get("topic", "").strip()
        count = int(body.get("slide_count", 5))
        user_id = body.get("user_id", "student")

        if not topic:
            return jsonify({"error": "Topic required"}), 400

        book_text, _, _ = get_memory(user_id)
        material = excerpt_search(book_text, topic, ctx_chars=2000) if book_text else topic

        system_prompt = (
            f"You are an expert presentation designer powered by Gemini AI. "
            f"Generate exactly {count} presentation slides for the topic: '{topic}'.\n"
            f"Output ONLY a valid JSON array of objects, where each object has keys: "
            f"'title', 'subtitle', and 'bullets' (a list of 3 bullet point strings).\n"
            f"Example: [{{\"title\": \"...\", \"subtitle\": \"...\", \"bullets\": [\"...\", \"...\"]}}]"
        )
        user_prompt = f"Study Material / Topic:\n{material}"

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        raw_llm = call_gemini(messages, max_tokens=1000, temperature=0.2, response_json=True)
        json_match = re.search(r'\[.*\]', raw_llm, re.DOTALL)
        if json_match:
            try:
                slides = json.loads(json_match.group(0))
                formatted_slides = []
                for i, s in enumerate(slides, 1):
                    formatted_slides.append({
                        "id": i,
                        "title": s.get("title") or f"Slide {i}: Key Concepts of {topic}",
                        "subtitle": s.get("subtitle") or f"Subtopic Analysis & Review {i}",
                        "bullets": s.get("bullets") or [f"Core principle {i}.1", f"Practical application {i}.2", f"Key takeaway {i}.3"]
                    })
                return jsonify({"success": True, "slides": formatted_slides})
            except Exception:
                pass

        # Fallback
        fallback_slides = [
            {
                "id": i + 1,
                "title": f"Slide {i + 1}: Key Aspects of {topic}",
                "subtitle": f"Overview & Analysis {i + 1}",
                "bullets": [
                    f"Primary concept definition and core principle {i + 1}.1",
                    f"Practical application and real-world example {i + 1}.2",
                    f"Key takeaway summary for review {i + 1}.3"
                ]
            } for i in range(count)
        ]
        return jsonify({"success": True, "slides": fallback_slides})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/export_ppt", methods=["POST"])
def export_ppt():
    try:
        body = request.get_json() or {}
        topic = body.get("topic", "Educational Presentation").strip()
        slides = body.get("slides", [])
        color_scheme = body.get("color_scheme", "indigo")
        template_style = body.get("template_style", "modern_cards")
        include_images = body.get("include_images", True)
        image_style = body.get("image_style", "minimalist")
        slide_shape = body.get("slide_shape", "rounded_card")
        slide_animation = body.get("slide_animation", "fade_glow")

        if not slides:
            return jsonify({"error": "No slides provided to export"}), 400

        filename = build_pptx_presentation(
            topic=topic,
            slides=slides,
            color_scheme=color_scheme,
            template_style=template_style,
            include_images=include_images,
            image_style=image_style,
            slide_shape=slide_shape,
            slide_animation=slide_animation
        )

        return jsonify({
            "success": True,
            "filename": filename,
            "download_url": f"/download/{filename}"
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate_flowchart", methods=["POST"])
def generate_flowchart():
    try:
        body = request.get_json() or {}
        topic = body.get("topic", "").strip()
        direction = body.get("direction", "TD")
        user_id = body.get("user_id", "student")

        if not topic:
            return jsonify({"error": "Topic required"}), 400

        book_text, _, _ = get_memory(user_id)
        material = excerpt_search(book_text, topic, ctx_chars=2000) if book_text else topic

        system_prompt = (
            f"You are a professional flowchart and system diagram architect powered by Gemini AI.\n"
            f"Generate a professional, high-quality, valid Mermaid.js flowchart for topic: '{topic}'.\n\n"
            f"Mermaid Syntax Rules:\n"
            f"1. Must start with 'flowchart {direction}' on the first line.\n"
            f"2. Node IDs must be alphanumeric and short (e.g., A, B, C, Node1, Node2). NEVER use spaces or special characters in node IDs.\n"
            f"3. Node shapes must use correct syntax:\n"
            f"   - Start/End: `Start([Start Text])` or `End([End Text])`\n"
            f"   - Process Box: `A[Process Text]`\n"
            f"   - Decision: `B{{Decision Question?}}`\n"
            f"   - Subprocess: `C[[Subprocess Text]]`\n"
            f"4. Links must connect node IDs directly. Connections must be on separate lines. Example:\n"
            f"   A --> B\n"
            f"   B -->|Yes| C\n"
            f"   B -->|No| D\n"
            f"5. NEVER nest node definitions or write multiple nodes/arrows inside brackets.\n"
            f"6. DO NOT use C-style comments like `//` or `/*`. Only use Mermaid comments starting with `%%` if needed, or avoid comments entirely.\n"
            f"7. Return ONLY the raw Mermaid syntax inside a ```mermaid ... ``` codeblock.\n\n"
            f"Example of Valid Mermaid Flowchart:\n"
            f"```mermaid\n"
            f"flowchart TD\n"
            f"    Start([Start Process]) --> Step1[Gather Requirements]\n"
            f"    Step1 --> Dec1{{Are requirements clear?}}\n"
            f"    Dec1 -->|Yes| Step2[Develop System Design]\n"
            f"    Dec1 -->|No| Step1\n"
            f"    Step2 --> End([Process Completed])\n"
            f"```"
        )
        user_prompt = f"Topic / Concept Details:\n{material}"

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        raw_llm = call_gemini(messages, max_tokens=600, temperature=0.2)
        
        mermaid_code = ""
        if "flowchart" in raw_llm:
            idx = raw_llm.find("flowchart")
            mermaid_code = raw_llm[idx:]
            if "```" in mermaid_code:
                mermaid_code = mermaid_code.split("```")[0].strip()
        
        if mermaid_code:
            # Post-process and sanitize common syntax mistakes made by LLM
            lines = []
            for line in mermaid_code.splitlines():
                stripped = line.strip()
                # Remove JS style comments
                if stripped.startswith("//") or stripped.startswith("/*"):
                    continue
                if "//" in line:
                    line = line.split("//")[0]
                lines.append(line)
            mermaid_code = "\n".join(lines).strip()
        
        if not mermaid_code:
            mermaid_code = (
                f"flowchart {direction}\n"
                f"    A([Start: {topic}]) --> B[Data Ingestion & Cleaning]\n"
                f"    B --> C{{Is Data Valid & Clean?}}\n"
                f"    C -->|No| D[Preprocessing & Outlier Removal]\n"
                f"    D --> B\n"
                f"    C -->|Yes| E[Feature Extraction & Engineering]\n"
                f"    E --> F[[Model Training & Cross-Validation]]\n"
                f"    F --> G{{Is Metric > Threshold?}}\n"
                f"    G -->|No| H[Hyperparameter Tuning]\n"
                f"    H --> F\n"
                f"    G -->|Yes| I[Model Deployment & Monitoring]\n"
                f"    I --> J([End: Pipeline Complete])"
            )

        return jsonify({
            "success": True,
            "mermaid_code": mermaid_code
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/socratic_hint", methods=["POST"])
def socratic_hint():
    try:
        body = request.get_json() or {}
        question = body.get("question", "").strip()
        user_answer = body.get("user_answer", "").strip()
        user_id = body.get("user_id", "student")

        if not question:
            return jsonify({"error": "Question empty"}), 400

        book_text, _, _ = get_memory(user_id)
        material = excerpt_search(book_text, question, ctx_chars=2000) if book_text else question

        system_prompt = (
            "You are a master Socratic AI Tutor powered by Gemini AI. "
            "Your goal is to guide the student to master the concept through Socratic inquiry.\n"
            "DO NOT just give away the answer! Instead:\n"
            "1. Evaluate their response gently (praise effort, highlight what was accurate).\n"
            "2. Provide a subtle Socratic Hint derived from the study material.\n"
            "3. Ask 1 targeted Counter-Question that pushes them to correct their mistake or think deeper.\n"
            "Keep your response concise, structured, and encouraging."
        )
        user_prompt = (
            f"Study Material Context:\n{material}\n\n"
            f"Question Asked: {question}\n"
            f"Student's Answer: {user_answer if user_answer else '[Student requested a hint]'}"
        )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        tutor_response = call_gemini(messages, max_tokens=350, temperature=0.3)
        return jsonify({
            "success": True,
            "socratic_response": tutor_response
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download/<filename>")
def download(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

if __name__ == "__main__":
    init_db()
    try:
        print("Pre-loading sentence-transformers embedding model...")
        get_embedding_model()
    except Exception as e:
        print("Warning: Could not pre-load embedding model:", e)
    print("Gemini AI Model initialized:", GEMINI_MODEL)
    app.run(debug=True, use_reloader=False, threaded=True)

